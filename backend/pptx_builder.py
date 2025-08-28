import os
import logging
from pathlib import Path
from typing import Dict, List, Any, Optional
import asyncio

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)


# backend/pptx_builder.py

from pptx import Presentation

def create_presentation_from_template(template_path: str, slides: list, output_path: str) -> str:
    """
    Creates a PowerPoint file from a template and a list of slides.
    slides = [{"title": "Slide 1", "content": "Some text"}, ...]
    """
    prs = Presentation(template_path)

    for slide_data in slides:
        slide_layout = prs.slide_layouts[1]  # Title & Content layout
        slide = prs.slides.add_slide(slide_layout)

        title_placeholder = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        title_placeholder.text = slide_data.get("title", "")
        content_placeholder.text = slide_data.get("content", "")

    prs.save(output_path)
    return output_path

async def create_presentation_from_template(
    template_path: str,
    output_path: str,
    presentation_data: Dict[str, Any],
    llm_client=None
) -> bool:
    """
    Create a PowerPoint presentation from a template and structured data
    
    Args:
        template_path: Path to the PowerPoint template file
        output_path: Path where the generated presentation will be saved
        presentation_data: Dictionary containing presentation structure
        llm_client: LLM client for generating additional content
    
    Returns:
        bool: True if successful, False otherwise
    """
    try:
        logger.info("🎨 Starting presentation creation")
        
        # Load template presentation
        prs = Presentation(template_path)
        logger.info(f"📄 Loaded template with {len(prs.slides)} slides")
        
        # Extract template information
        template_info = extract_template_info(prs)
        logger.info(f"🎨 Template info: {len(template_info['layouts'])} layouts, {len(template_info['colors'])} colors")
        
        # Clear existing slides (keep layouts)
        slide_count = len(prs.slides)
        for i in range(slide_count - 1, -1, -1):
            r_id = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(r_id)
            del prs.slides._sldIdLst[i]
        
        # Create slides from data
        slides_data = presentation_data.get('slides', [])
        if not slides_data:
            logger.warning("⚠️ No slides data provided")
            return False
        
        logger.info(f"📊 Creating {len(slides_data)} slides")
        
        # Create title slide
        title_slide_data = {
            'title': presentation_data.get('title', 'AI Generated Presentation'),
            'content': presentation_data.get('subtitle', ''),
            'type': 'title'
        }
        await create_slide(prs, title_slide_data, template_info, 0)
        
        # Create content slides
        for i, slide_data in enumerate(slides_data, 1):
            await create_slide(prs, slide_data, template_info, i)
        
        # Save presentation
        prs.save(output_path)
        logger.info(f"✅ Presentation saved to {output_path}")
        
        return True
        
    except Exception as e:
        logger.error(f"❌ Error creating presentation: {e}", exc_info=True)
        return False

def extract_template_info(prs: Presentation) -> Dict[str, Any]:
    """Extract styling information from the template"""
    try:
        template_info = {
            'layouts': [],
            'colors': [],
            'fonts': [],
            'images': []
        }
        
        # Extract slide layouts
        for layout in prs.slide_layouts:
            layout_info = {
                'name': layout.name if hasattr(layout, 'name') else f"Layout {len(template_info['layouts'])}",
                'placeholders': []
            }
            
            # Extract placeholders
            for placeholder in layout.placeholders:
                try:
                    ph_info = {
                        'type': placeholder.placeholder_format.type if hasattr(placeholder, 'placeholder_format') else None,
                        'idx': placeholder.placeholder_format.idx if hasattr(placeholder, 'placeholder_format') else len(layout_info['placeholders']),
                        'left': placeholder.left,
                        'top': placeholder.top,
                        'width': placeholder.width,
                        'height': placeholder.height
                    }
                    layout_info['placeholders'].append(ph_info)
                except:
                    pass
            
            template_info['layouts'].append(layout_info)
        
        # Extract color scheme from existing slides
        for slide in prs.slides:
            for shape in slide.shapes:
                try:
                    if hasattr(shape, 'fill') and shape.fill.type == 1:  # Solid fill
                        color = shape.fill.fore_color.rgb
                        if color not in template_info['colors']:
                            template_info['colors'].append(color)
                except:
                    pass
        
        # Add default colors if none found
        if not template_info['colors']:
            template_info['colors'] = [
                RGBColor(44, 62, 80),   # Dark blue
                RGBColor(52, 152, 219), # Blue
                RGBColor(46, 204, 113), # Green
                RGBColor(241, 196, 15), # Yellow
                RGBColor(231, 76, 60),  # Red
                RGBColor(155, 89, 182), # Purple
                RGBColor(149, 165, 166) # Gray
            ]
        
        return template_info
        
    except Exception as e:
        logger.warning(f"⚠️ Error extracting template info: {e}")
        return {
            'layouts': [],
            'colors': [RGBColor(44, 62, 80), RGBColor(52, 152, 219)],
            'fonts': ['Calibri', 'Arial'],
            'images': []
        }

async def create_slide(prs: Presentation, slide_data: Dict[str, Any], template_info: Dict[str, Any], slide_index: int):
    """Create a single slide from data"""
    try:
        slide_type = slide_data.get('type', 'content')
        
        # Choose appropriate layout
        layout_index = 0  # Title slide
        if slide_type in ['content', 'bullets']:
            layout_index = min(1, len(prs.slide_layouts) - 1)  # Content layout
        elif slide_type == 'title':
            layout_index = 0  # Title layout
        elif slide_type in ['image', 'chart']:
            layout_index = min(2, len(prs.slide_layouts) - 1) if len(prs.slide_layouts) > 2 else 1
        
        layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(layout)
        
        logger.info(f"📄 Creating slide {slide_index + 1}: {slide_data.get('title', 'Untitled')}")
        
        # Add title
        title = slide_data.get('title', f'Slide {slide_index + 1}')
        if slide.shapes.title:
            slide.shapes.title.text = title
            
            # Style title
            title_frame = slide.shapes.title.text_frame
            title_frame.clear()
            p = title_frame.paragraphs[0]
            p.text = title
            p.font.name = 'Calibri'
            p.font.size = Pt(36 if slide_type == 'title' else 28)
            p.font.bold = True
            if template_info['colors']:
                p.font.color.rgb = template_info['colors'][0]
        
        # Add content based on slide type
        if slide_type == 'title':
            await add_title_content(slide, slide_data, template_info)
        elif slide_type == 'bullets':
            await add_bullet_content(slide, slide_data, template_info)
        elif slide_type == 'content':
            await add_content_slide(slide, slide_data, template_info)
        else:
            await add_content_slide(slide, slide_data, template_info)
        
        # Add speaker notes
        notes_text = slide_data.get('speaker_notes', '')
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text
            
    except Exception as e:
        logger.error(f"❌ Error creating slide {slide_index}: {e}")

async def add_title_content(slide, slide_data: Dict[str, Any], template_info: Dict[str, Any]):
    """Add content to title slide"""
    try:
        content = slide_data.get('content', '')
        
        # Look for subtitle placeholder or content placeholder
        subtitle_shape = None
        for shape in slide.shapes:
            if hasattr(shape, 'placeholder_format'):
                # Subtitle placeholder
                if shape.placeholder_format.type == 4:  # PP_PLACEHOLDER.SUBTITLE
                    subtitle_shape = shape
                    break
        
        if subtitle_shape and content:
            subtitle_shape.text = content
            
            # Style subtitle
            text_frame = subtitle_shape.text_frame
            for paragraph in text_frame.paragraphs:
                paragraph.font.name = 'Calibri'
                paragraph.font.size = Pt(18)
                if template_info['colors'] and len(template_info['colors']) > 1:
                    paragraph.font.color.rgb = template_info['colors'][1]
                    
    except Exception as e:
        logger.error(f"❌ Error adding title content: {e}")

async def add_bullet_content(slide, slide_data: Dict[str, Any], template_info: Dict[str, Any]):
    """Add bullet point content to slide"""
    try:
        content = slide_data.get('content', [])
        if isinstance(content, str):
            content = [content]
        
        # Find content placeholder
        content_shape = None
        for shape in slide.shapes:
            if hasattr(shape, 'placeholder_format'):
                if shape.placeholder_format.type == 2:  # PP_PLACEHOLDER.BODY
                    content_shape = shape
                    break
        
        if content_shape and content:
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            for i, bullet_point in enumerate(content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = str(bullet_point).strip()
                p.level = 0
                p.font.name = 'Calibri'
                p.font.size = Pt(18)
                
                if template_info['colors']:
                    p.font.color.rgb = template_info['colors'][0]
                    
    except Exception as e:
        logger.error(f"❌ Error adding bullet content: {e}")

async def add_content_slide(slide, slide_data: Dict[str, Any], template_info: Dict[str, Any]):
    """Add general content to slide"""
    try:
        content = slide_data.get('content', '')
        
        # Find content placeholder
        content_shape = None
        for shape in slide.shapes:
            if hasattr(shape, 'placeholder_format'):
                if shape.placeholder_format.type == 2:  # PP_PLACEHOLDER.BODY
                    content_shape = shape
                    break
        
        if not content_shape:
            # Create a text box if no placeholder found
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4)
            content_shape = slide.shapes.add_textbox(left, top, width, height)
        
        if content and content_shape:
            if isinstance(content, list):
                # Handle list content as bullet points
                text_frame = content_shape.text_frame
                text_frame.clear()
                
                for i, item in enumerate(content):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = str(item).strip()
                    p.level = 0
                    p.font.name = 'Calibri'
                    p.font.size = Pt(16)
                    
                    if template_info['colors']:
                        p.font.color.rgb = template_info['colors'][0]
            else:
                # Handle string content
                content_shape.text = str(content)
                text_frame = content_shape.text_frame
                
                for paragraph in text_frame.paragraphs:
                    paragraph.font.name = 'Calibri'
                    paragraph.font.size = Pt(16)
                    
                    if template_info['colors']:
                        paragraph.font.color.rgb = template_info['colors'][0]
                        
    except Exception as e:
        logger.error(f"❌ Error adding content: {e}")

def validate_presentation_data(data: Dict[str, Any]) -> bool:
    """Validate presentation data structure"""
    try:
        if not isinstance(data, dict):
            return False
        
        # Check required fields
        if 'slides' not in data or not isinstance(data['slides'], list):
            return False
        
        if len(data['slides']) == 0:
            return False
        
        # Check each slide
        for slide in data['slides']:
            if not isinstance(slide, dict):
                return False
            
            if 'title' not in slide:
                return False
        
        return True
        
    except Exception:
        return False
